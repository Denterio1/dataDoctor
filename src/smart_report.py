"""
dataDoctor — src/smart_report.py
==================================
Professional Multi-Format Report Generator

Formats:
    ▸ PDF   — ReportLab: cover page, charts, tables, watermark, header/footer
    ▸ DOCX  — python-docx: styled document with charts via matplotlib
    ▸ PPTX  — python-pptx: full presentation with charts + speaker notes

All formats receive:
    ▸ Executive Summary
    ▸ Overall Quality Score + Grade
    ▸ 7-Dimension breakdown with bar chart
    ▸ Issues & Recommendations table
    ▸ Dataset Statistics table
    ▸ Per-column null heatmap (PDF/DOCX)
    ▸ PII alert section
    ▸ Branding: dataDoctor header/footer

Author  : Kader (Denterio1)
Version : 1.0.0
"""

from __future__ import annotations

import io
import os
import logging
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch

logger = logging.getLogger("dataDoctor.smart_report")

# ── Brand colors ─────────────────────────────────────────────────────────────
BRAND_PRIMARY   = "#1E3A5F"    # deep navy
BRAND_ACCENT    = "#00B4D8"    # cyan
BRAND_SUCCESS   = "#2DC653"    # green
BRAND_WARNING   = "#F4A261"    # orange
BRAND_DANGER    = "#E63946"    # red
BRAND_LIGHT     = "#F8F9FA"    # almost white
BRAND_DARK      = "#212529"    # near black
BRAND_GRAY      = "#6C757D"    # medium gray

GRADE_COLORS = {
    "A": BRAND_SUCCESS,
    "B": "#74C69D",
    "C": BRAND_WARNING,
    "D": "#E76F51",
    "F": BRAND_DANGER,
}

SCORE_COLOR_MAP = [
    (0,   40,  BRAND_DANGER),
    (40,  60,  "#E76F51"),
    (60,  75,  BRAND_WARNING),
    (75,  90,  "#74C69D"),
    (90,  101, BRAND_SUCCESS),
]


def _score_color(score: float) -> str:
    for lo, hi, color in SCORE_COLOR_MAP:
        if lo <= score < hi:
            return color
    return BRAND_SUCCESS


# ─────────────────────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
#  CHART FACTORY — matplotlib charts saved to bytes/files
# ══════════════════════════════════════════════════════════════════════════════
# ─────────────────────────────────────────────────────────────────────────────

class ChartFactory:
    """Generates all matplotlib charts used across PDF, DOCX, PPTX."""

    DPI = 150

    # ── dimension bar chart ───────────────────────────────────────────────────
    @staticmethod
    def dimension_bar_chart(
        dimensions: Dict[str, Any],
        overall: float,
        title: str = "Data Quality — Dimension Scores",
        figsize: Tuple[float, float] = (9, 5),
    ) -> bytes:
        names  = [d.name for d in dimensions.values()]
        scores = [d.score for d in dimensions.values()]
        colors = [_score_color(s) for s in scores]

        fig, ax = plt.subplots(figsize=figsize, facecolor=BRAND_LIGHT)
        ax.set_facecolor(BRAND_LIGHT)

        bars = ax.barh(names, scores, color=colors, height=0.6,
                       edgecolor="white", linewidth=0.8)

        # value labels
        for bar, score in zip(bars, scores):
            ax.text(
                min(score + 1.5, 97), bar.get_y() + bar.get_height() / 2,
                f"{score:.1f}",
                va="center", ha="left", fontsize=10,
                color=BRAND_DARK, fontweight="bold",
            )

        # overall line
        ax.axvline(overall, color=BRAND_PRIMARY, linewidth=2,
                   linestyle="--", alpha=0.7, label=f"Overall: {overall:.1f}")

        ax.set_xlim(0, 110)
        ax.set_xlabel("Score (0–100)", fontsize=10, color=BRAND_GRAY)
        ax.set_title(title, fontsize=13, fontweight="bold",
                     color=BRAND_PRIMARY, pad=12)
        ax.legend(fontsize=9, loc="lower right")
        ax.spines[["top", "right"]].set_visible(False)
        ax.tick_params(colors=BRAND_DARK)

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=ChartFactory.DPI,
                    bbox_inches="tight", facecolor=BRAND_LIGHT)
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    # ── score gauge / donut chart ─────────────────────────────────────────────
    @staticmethod
    def score_gauge(
        score: float,
        grade: str,
        figsize: Tuple[float, float] = (4, 4),
    ) -> bytes:
        fig, ax = plt.subplots(figsize=figsize, facecolor="white",
                               subplot_kw=dict(aspect="equal"))
        ax.set_facecolor("white")

        # background ring
        theta1, theta2 = 0, 360
        ring_bg = mpatches.Wedge(
            (0.5, 0.5), 0.42, theta1, theta2,
            width=0.15, transform=ax.transAxes,
            facecolor="#E9ECEF", edgecolor="none",
        )
        ax.add_patch(ring_bg)

        # filled arc based on score
        arc_deg = score / 100 * 360
        ring_fg = mpatches.Wedge(
            (0.5, 0.5), 0.42, 90, 90 - arc_deg,
            width=0.15, transform=ax.transAxes,
            facecolor=_score_color(score),
            edgecolor="none",
        )
        ax.add_patch(ring_fg)

        # center text
        ax.text(0.5, 0.55, f"{score:.0f}", ha="center", va="center",
                fontsize=28, fontweight="bold", color=BRAND_DARK,
                transform=ax.transAxes)
        ax.text(0.5, 0.38, f"Grade {grade}", ha="center", va="center",
                fontsize=12, color=GRADE_COLORS.get(grade, BRAND_GRAY),
                fontweight="bold", transform=ax.transAxes)
        ax.text(0.5, 0.28, "/ 100", ha="center", va="center",
                fontsize=10, color=BRAND_GRAY, transform=ax.transAxes)

        ax.axis("off")
        ax.set_title("Overall Score", fontsize=11, color=BRAND_PRIMARY,
                     fontweight="bold", pad=4)

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=ChartFactory.DPI,
                    bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    # ── null heatmap ──────────────────────────────────────────────────────────
    @staticmethod
    def null_heatmap(df: pd.DataFrame, figsize: Tuple[float, float] = (10, 4)) -> bytes:
        null_rates = df.isnull().mean().sort_values(ascending=False)
        null_rates = null_rates[null_rates > 0]

        if len(null_rates) == 0:
            # all clean — simple text chart
            fig, ax = plt.subplots(figsize=(5, 2), facecolor="white")
            ax.text(0.5, 0.5, "✅  No missing values found",
                    ha="center", va="center", fontsize=14,
                    color=BRAND_SUCCESS, fontweight="bold",
                    transform=ax.transAxes)
            ax.axis("off")
        else:
            fig, ax = plt.subplots(figsize=figsize, facecolor=BRAND_LIGHT)
            ax.set_facecolor(BRAND_LIGHT)
            cols   = null_rates.index.tolist()
            values = null_rates.values

            bar_colors = [_score_color(100 - v * 100) for v in values]
            bars = ax.bar(range(len(cols)), values * 100,
                          color=bar_colors, edgecolor="white", linewidth=0.5)

            for bar, val in zip(bars, values):
                ax.text(bar.get_x() + bar.get_width() / 2,
                        bar.get_height() + 0.5,
                        f"{val:.1%}", ha="center", va="bottom",
                        fontsize=8, color=BRAND_DARK)

            ax.set_xticks(range(len(cols)))
            ax.set_xticklabels(cols, rotation=45, ha="right", fontsize=9)
            ax.set_ylabel("Missing Rate (%)", fontsize=10)
            ax.set_title("Missing Values by Column", fontsize=12,
                         fontweight="bold", color=BRAND_PRIMARY)
            ax.spines[["top", "right"]].set_visible(False)
            ax.set_ylim(0, min(105, values.max() * 100 * 1.2))

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=ChartFactory.DPI,
                    bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    # ── distribution grid ─────────────────────────────────────────────────────
    @staticmethod
    def distribution_grid(
        df: pd.DataFrame,
        max_cols: int = 6,
        figsize: Tuple[float, float] = (12, 6),
    ) -> bytes:
        num_cols = df.select_dtypes(include=[np.number]).columns.tolist()[:max_cols]
        if not num_cols:
            fig, ax = plt.subplots(figsize=(4, 2))
            ax.text(0.5, 0.5, "No numeric columns", ha="center", va="center")
            ax.axis("off")
        else:
            n    = len(num_cols)
            cols = min(3, n)
            rows = (n + cols - 1) // cols
            fig, axes = plt.subplots(rows, cols,
                                     figsize=(figsize[0], figsize[1] * rows / 2),
                                     facecolor=BRAND_LIGHT)
            axes = np.array(axes).flatten()

            for i, col in enumerate(num_cols):
                s = df[col].dropna()
                axes[i].hist(s, bins=30, color=BRAND_ACCENT,
                             edgecolor="white", alpha=0.85)
                axes[i].set_title(col, fontsize=9, color=BRAND_PRIMARY,
                                  fontweight="bold")
                axes[i].spines[["top", "right"]].set_visible(False)
                axes[i].tick_params(labelsize=7)

            for j in range(i + 1, len(axes)):
                axes[j].axis("off")

            fig.suptitle("Numeric Column Distributions",
                         fontsize=12, fontweight="bold", color=BRAND_PRIMARY)

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=ChartFactory.DPI, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    # ── recommendations chart ─────────────────────────────────────────────────
    @staticmethod
    def issue_priority_chart(
        dimensions: Dict[str, Any],
        figsize: Tuple[float, float] = (8, 4),
    ) -> bytes:
        names  = [d.name for d in dimensions.values()]
        issues = [len(d.issues) for d in dimensions.values()]
        colors = [_score_color(d.score) for d in dimensions.values()]

        fig, ax = plt.subplots(figsize=figsize, facecolor=BRAND_LIGHT)
        ax.set_facecolor(BRAND_LIGHT)

        bars = ax.bar(names, issues, color=colors, edgecolor="white",
                      linewidth=0.8, width=0.6)

        for bar, cnt in zip(bars, issues):
            if cnt > 0:
                ax.text(bar.get_x() + bar.get_width() / 2,
                        bar.get_height() + 0.05,
                        str(cnt), ha="center", va="bottom",
                        fontsize=10, fontweight="bold", color=BRAND_DARK)

        ax.set_ylabel("Issue Count", fontsize=10, color=BRAND_GRAY)
        ax.set_title("Issues by Dimension", fontsize=12,
                     fontweight="bold", color=BRAND_PRIMARY)
        ax.spines[["top", "right"]].set_visible(False)
        ax.tick_params(axis="x", rotation=20)

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=ChartFactory.DPI, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf.read()


# ─────────────────────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
#  PDF REPORT GENERATOR
# ══════════════════════════════════════════════════════════════════════════════
# ─────────────────────────────────────────────────────────────────────────────

class PDFReportGenerator:
    """
    Generates a professional multi-page PDF report using ReportLab.

    Pages:
        1. Cover page
        2. Executive Summary + Score Gauge
        3. Dimension Breakdown + Bar Chart
        4. Issues & Recommendations
        5. Dataset Statistics
        6. Missing Values Heatmap
        7. Distributions
        8. PII Alert (if any)
    """

    def __init__(self, brand_name: str = "dataDoctor"):
        self.brand = brand_name

    def _rl_color(self, hex_color: str):
        from reportlab.lib import colors as rlcolors
        h = hex_color.lstrip("#")
        r, g, b = tuple(int(h[i:i+2], 16) / 255.0 for i in (0, 2, 4))
        return rlcolors.Color(r, g, b)

    def _build_styles(self):
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
        from reportlab.lib import colors as rlcolors

        styles = getSampleStyleSheet()
        primary = self._rl_color(BRAND_PRIMARY)

        styles.add(ParagraphStyle(
            "DDTitle",
            fontName="Helvetica-Bold", fontSize=28,
            textColor=primary, alignment=TA_CENTER, spaceAfter=8,
        ))
        styles.add(ParagraphStyle(
            "DDSubtitle",
            fontName="Helvetica", fontSize=14,
            textColor=self._rl_color(BRAND_GRAY),
            alignment=TA_CENTER, spaceAfter=4,
        ))
        styles.add(ParagraphStyle(
            "DDH1",
            fontName="Helvetica-Bold", fontSize=16,
            textColor=primary, spaceBefore=14, spaceAfter=6,
        ))
        styles.add(ParagraphStyle(
            "DDH2",
            fontName="Helvetica-Bold", fontSize=13,
            textColor=primary, spaceBefore=10, spaceAfter=4,
        ))
        styles.add(ParagraphStyle(
            "DDBody",
            fontName="Helvetica", fontSize=10,
            textColor=self._rl_color(BRAND_DARK),
            leading=14, spaceAfter=4,
        ))
        styles.add(ParagraphStyle(
            "DDSmall",
            fontName="Helvetica", fontSize=8,
            textColor=self._rl_color(BRAND_GRAY),
            leading=11,
        ))
        styles.add(ParagraphStyle(
            "DDAlert",
            fontName="Helvetica-Bold", fontSize=10,
            textColor=self._rl_color(BRAND_DANGER),
            leading=14, spaceAfter=4,
        ))
        return styles

    def _header_footer(self, canvas, doc):
        from reportlab.lib.units import mm
        from reportlab.lib import colors as rlcolors

        canvas.saveState()
        w, h = doc.pagesize

        # header bar
        canvas.setFillColor(self._rl_color(BRAND_PRIMARY))
        canvas.rect(0, h - 18 * mm, w, 18 * mm, fill=1, stroke=0)
        canvas.setFillColor(rlcolors.white)
        canvas.setFont("Helvetica-Bold", 11)
        canvas.drawString(15 * mm, h - 11 * mm, f"🩺 {self.brand}")
        canvas.setFont("Helvetica", 9)
        canvas.drawRightString(w - 15 * mm, h - 11 * mm,
                               "Data Quality Intelligence Report")

        # footer
        canvas.setFillColor(self._rl_color(BRAND_GRAY))
        canvas.rect(0, 0, w, 10 * mm, fill=1, stroke=0)
        canvas.setFillColor(rlcolors.white)
        canvas.setFont("Helvetica", 8)
        canvas.drawString(15 * mm, 3.5 * mm,
                          f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        canvas.drawCentredString(w / 2, 3.5 * mm,
                                 "Confidential — dataDoctor Report")
        canvas.drawRightString(w - 15 * mm, 3.5 * mm,
                               f"Page {doc.page}")

        canvas.restoreState()

    def _img_flowable(self, img_bytes: bytes, width_mm: float):
        from reportlab.platypus import Image as RLImage
        from reportlab.lib.units import mm
        buf = io.BytesIO(img_bytes)
        img = RLImage(buf)
        aspect = img.imageHeight / float(img.imageWidth)
        w = width_mm * mm
        return RLImage(io.BytesIO(img_bytes), width=w, height=w * aspect)

    def _dimension_table(self, dimensions: Dict[str, Any], styles):
        from reportlab.platypus import Table, TableStyle
        from reportlab.lib.units import mm
        from reportlab.lib import colors as rlcolors

        data = [["Dimension", "Score", "Weight", "Status", "Issues"]]
        for dim in dimensions.values():
            status = (
                "✅ Excellent" if dim.score >= 85 else
                "🟡 Good"     if dim.score >= 70 else
                "🟠 Fair"     if dim.score >= 55 else
                "🔴 Poor"
            )
            data.append([
                dim.name,
                f"{dim.score:.1f}/100",
                f"{dim.weight * 100:.0f}%",
                status,
                str(len(dim.issues)),
            ])

        col_widths = [42*mm, 22*mm, 18*mm, 32*mm, 16*mm]
        tbl = Table(data, colWidths=col_widths, repeatRows=1)
        tbl.setStyle(TableStyle([
            # header
            ("BACKGROUND",   (0, 0), (-1, 0), self._rl_color(BRAND_PRIMARY)),
            ("TEXTCOLOR",    (0, 0), (-1, 0), rlcolors.white),
            ("FONTNAME",     (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",     (0, 0), (-1, 0), 9),
            ("ALIGN",        (0, 0), (-1, 0), "CENTER"),
            ("BOTTOMPADDING",(0, 0), (-1, 0), 7),
            ("TOPPADDING",   (0, 0), (-1, 0), 7),
            # body
            ("FONTNAME",     (0, 1), (-1, -1), "Helvetica"),
            ("FONTSIZE",     (0, 1), (-1, -1), 9),
            ("ROWBACKGROUNDS",(0, 1), (-1, -1),
             [self._rl_color(BRAND_LIGHT), rlcolors.white]),
            ("ALIGN",        (1, 1), (-1, -1), "CENTER"),
            ("VALIGN",       (0, 0), (-1, -1), "MIDDLE"),
            ("BOTTOMPADDING",(0, 1), (-1, -1), 5),
            ("TOPPADDING",   (0, 1), (-1, -1), 5),
            ("GRID",         (0, 0), (-1, -1), 0.4, self._rl_color(BRAND_GRAY)),
        ]))
        return tbl

    def _stats_table(self, stats: Dict[str, Any], styles):
        from reportlab.platypus import Table, TableStyle
        from reportlab.lib.units import mm
        from reportlab.lib import colors as rlcolors

        data = [["Metric", "Value"]]
        labels = {
            "rows"             : "Total Rows",
            "columns"          : "Total Columns",
            "numeric_columns"  : "Numeric Columns",
            "categorical_cols" : "Categorical Columns",
            "datetime_cols"    : "Datetime Columns",
            "null_cells"       : "Null Cells",
            "null_rate"        : "Overall Null Rate",
            "duplicate_rows"   : "Duplicate Rows",
            "memory_mb"        : "Memory Usage (MB)",
        }
        for key, label in labels.items():
            val = stats.get(key, "N/A")
            if key == "null_rate":
                val = f"{float(val):.2%}"
            data.append([label, str(val)])

        col_widths = [70*mm, 50*mm]
        tbl = Table(data, colWidths=col_widths, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND",   (0, 0), (-1, 0), self._rl_color(BRAND_PRIMARY)),
            ("TEXTCOLOR",    (0, 0), (-1, 0), rlcolors.white),
            ("FONTNAME",     (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",     (0, 0), (-1, 0), 9),
            ("FONTNAME",     (0, 1), (-1, -1), "Helvetica"),
            ("FONTSIZE",     (0, 1), (-1, -1), 9),
            ("ROWBACKGROUNDS",(0, 1), (-1, -1),
             [self._rl_color(BRAND_LIGHT), rlcolors.white]),
            ("ALIGN",        (1, 1), (-1, -1), "CENTER"),
            ("GRID",         (0, 0), (-1, -1), 0.4, self._rl_color(BRAND_GRAY)),
            ("BOTTOMPADDING",(0, 0), (-1, -1), 5),
            ("TOPPADDING",   (0, 0), (-1, -1), 5),
        ]))
        return tbl

    def _recs_table(self, recommendations: List[str], styles):
        from reportlab.platypus import Table, TableStyle, Paragraph
        from reportlab.lib.units import mm
        from reportlab.lib import colors as rlcolors

        data = [["#", "Recommendation"]]
        for i, rec in enumerate(recommendations, 1):
            data.append([str(i), Paragraph(rec, styles["DDSmall"])])

        col_widths = [8*mm, 152*mm]
        tbl = Table(data, colWidths=col_widths, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND",   (0, 0), (-1, 0), self._rl_color(BRAND_ACCENT)),
            ("TEXTCOLOR",    (0, 0), (-1, 0), rlcolors.white),
            ("FONTNAME",     (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",     (0, 0), (-1, 0), 9),
            ("FONTNAME",     (0, 1), (-1, -1), "Helvetica"),
            ("FONTSIZE",     (0, 1), (-1, -1), 9),
            ("ROWBACKGROUNDS",(0, 1), (-1, -1),
             [self._rl_color(BRAND_LIGHT), rlcolors.white]),
            ("ALIGN",        (0, 1), (0, -1), "CENTER"),
            ("VALIGN",       (0, 0), (-1, -1), "MIDDLE"),
            ("GRID",         (0, 0), (-1, -1), 0.4, self._rl_color(BRAND_GRAY)),
            ("BOTTOMPADDING",(0, 0), (-1, -1), 5),
            ("TOPPADDING",   (0, 0), (-1, -1), 5),
        ]))
        return tbl

    def generate(
        self,
        profile,
        df: pd.DataFrame,
        output_path: str,
        dataset_name: str = "Dataset",
    ) -> str:
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer,
            PageBreak, HRFlowable, KeepTogether,
        )
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib import colors as rlcolors

        doc = SimpleDocTemplate(
            output_path,
            pagesize=A4,
            rightMargin=18*mm, leftMargin=18*mm,
            topMargin=24*mm, bottomMargin=18*mm,
            title=f"dataDoctor Quality Report — {dataset_name}",
            author="dataDoctor",
            subject="Data Quality Intelligence Report",
        )

        styles = self._build_styles()
        story  = []
        W      = A4[0] - 36*mm   # usable width

        # ── COVER PAGE ────────────────────────────────────────────────────────
        story.append(Spacer(1, 30*mm))
        story.append(Paragraph("🩺 dataDoctor", styles["DDTitle"]))
        story.append(Paragraph("Data Quality Intelligence Report", styles["DDSubtitle"]))
        story.append(Spacer(1, 6*mm))
        story.append(HRFlowable(
            width="100%", thickness=2,
            color=self._rl_color(BRAND_ACCENT), spaceAfter=6*mm,
        ))
        story.append(Paragraph(f"Dataset: <b>{dataset_name}</b>", styles["DDSubtitle"]))
        story.append(Paragraph(
            f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}",
            styles["DDSmall"],
        ))
        story.append(Spacer(1, 10*mm))

        # score gauge on cover
        gauge_bytes = ChartFactory.score_gauge(
            profile.overall_score, profile.grade, figsize=(4, 4)
        )
        story.append(self._img_flowable(gauge_bytes, 70))
        story.append(Spacer(1, 6*mm))

        grade_label = profile.grade_label
        story.append(Paragraph(
            f"<b>Overall Quality Score: {profile.overall_score:.1f}/100</b><br/>"
            f"{grade_label}  |  {profile.health_badge}",
            styles["DDH2"],
        ))
        story.append(Paragraph(
            f"Shape: {profile.dataset_shape[0]:,} rows × {profile.dataset_shape[1]} columns  |  "
            f"{len(profile.numeric_columns)} numeric features",
            styles["DDBody"],
        ))
        story.append(PageBreak())

        # ── EXECUTIVE SUMMARY ─────────────────────────────────────────────────
        story.append(Paragraph("1. Executive Summary", styles["DDH1"]))
        story.append(HRFlowable(width="100%", thickness=1,
                                color=self._rl_color(BRAND_ACCENT)))
        story.append(Spacer(1, 3*mm))

        total_issues = sum(len(d.issues) for d in profile.dimensions.values())
        total_pos    = sum(len(d.positives) for d in profile.dimensions.values())

        summary_text = (
            f"This report presents a comprehensive data quality analysis of the dataset "
            f"<b>{dataset_name}</b>, evaluated across 7 quality dimensions using the "
            f"dataDoctor Quality Scoring Engine. "
            f"The dataset contains <b>{profile.dataset_shape[0]:,} rows</b> and "
            f"<b>{profile.dataset_shape[1]} columns</b>. "
            f"<br/><br/>"
            f"The overall quality score is <b>{profile.overall_score:.1f}/100</b> "
            f"({grade_label}). "
            f"A total of <b>{total_issues} issues</b> were identified across all dimensions, "
            f"alongside <b>{total_pos} positive indicators</b>. "
            f"<br/><br/>"
        )
        if profile.pii_detected:
            summary_text += (
                f"<font color='red'><b>⚠ PII ALERT:</b> Personally identifiable information "
                f"was detected in {len(profile.pii_detected)} column(s). "
                f"Review the Integrity section for details.</font><br/><br/>"
            )
        summary_text += profile.recommendations[0] if profile.recommendations else ""
        story.append(Paragraph(summary_text, styles["DDBody"]))
        story.append(Spacer(1, 4*mm))

        # dimension bar chart
        story.append(Paragraph("Dimension Scores Overview", styles["DDH2"]))
        bar_bytes = ChartFactory.dimension_bar_chart(
            profile.dimensions, profile.overall_score, figsize=(9, 4.5)
        )
        story.append(self._img_flowable(bar_bytes, 160))
        story.append(PageBreak())

        # ── DIMENSION BREAKDOWN ───────────────────────────────────────────────
        story.append(Paragraph("2. Dimension Breakdown", styles["DDH1"]))
        story.append(HRFlowable(width="100%", thickness=1,
                                color=self._rl_color(BRAND_ACCENT)))
        story.append(Spacer(1, 3*mm))
        story.append(self._dimension_table(profile.dimensions, styles))
        story.append(Spacer(1, 6*mm))

        # per-dimension issues
        for dim in profile.dimensions.values():
            if not dim.issues and not dim.positives:
                continue
            story.append(Paragraph(f"2.{list(profile.dimensions.keys()).index(dim.name.lower().replace(' ', '_')) + 1}  {dim.name} — {dim.score:.1f}/100", styles["DDH2"]))
            if dim.positives:
                for p in dim.positives[:3]:
                    story.append(Paragraph(f"&nbsp;&nbsp;{p}", styles["DDBody"]))
            if dim.issues:
                for issue in dim.issues[:4]:
                    story.append(Paragraph(f"&nbsp;&nbsp;⚠ {issue}", styles["DDAlert"]))
            story.append(Spacer(1, 2*mm))

        story.append(PageBreak())

        # ── ISSUES & RECOMMENDATIONS ──────────────────────────────────────────
        story.append(Paragraph("3. Issues & Recommendations", styles["DDH1"]))
        story.append(HRFlowable(width="100%", thickness=1,
                                color=self._rl_color(BRAND_ACCENT)))
        story.append(Spacer(1, 3*mm))

        # issue priority chart
        issue_bytes = ChartFactory.issue_priority_chart(
            profile.dimensions, figsize=(8, 3.5)
        )
        story.append(self._img_flowable(issue_bytes, 140))
        story.append(Spacer(1, 4*mm))

        story.append(Paragraph("Prioritized Recommendations:", styles["DDH2"]))
        story.append(self._recs_table(profile.recommendations, styles))
        story.append(PageBreak())

        # ── DATASET STATISTICS ────────────────────────────────────────────────
        story.append(Paragraph("4. Dataset Statistics", styles["DDH1"]))
        story.append(HRFlowable(width="100%", thickness=1,
                                color=self._rl_color(BRAND_ACCENT)))
        story.append(Spacer(1, 3*mm))
        story.append(self._stats_table(profile.statistics, styles))
        story.append(Spacer(1, 6*mm))

        # column types breakdown
        story.append(Paragraph("Column Summary", styles["DDH2"]))
        col_data = [["Column", "Dtype", "Null %", "Unique", "Most Common"]]
        for col in df.columns[:20]:  # max 20
            null_pct = f"{df[col].isnull().mean():.1%}"
            unique   = str(df[col].nunique())
            try:
                most_common = str(df[col].value_counts().index[0])[:20]
            except Exception:
                most_common = "—"
            col_data.append([col[:25], str(df[col].dtype), null_pct, unique, most_common])

        from reportlab.platypus import Table, TableStyle
        from reportlab.lib import colors as rlcolors
        col_tbl = Table(col_data, colWidths=[45*mm, 25*mm, 18*mm, 20*mm, 50*mm],
                        repeatRows=1)
        col_tbl.setStyle(TableStyle([
            ("BACKGROUND",   (0, 0), (-1, 0), self._rl_color(BRAND_PRIMARY)),
            ("TEXTCOLOR",    (0, 0), (-1, 0), rlcolors.white),
            ("FONTNAME",     (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",     (0, 0), (-1, 0), 8),
            ("FONTNAME",     (0, 1), (-1, -1), "Helvetica"),
            ("FONTSIZE",     (0, 1), (-1, -1), 8),
            ("ROWBACKGROUNDS",(0, 1), (-1, -1),
             [self._rl_color(BRAND_LIGHT), rlcolors.white]),
            ("ALIGN",        (1, 1), (-1, -1), "CENTER"),
            ("GRID",         (0, 0), (-1, -1), 0.3, self._rl_color(BRAND_GRAY)),
            ("BOTTOMPADDING",(0, 0), (-1, -1), 4),
            ("TOPPADDING",   (0, 0), (-1, -1), 4),
        ]))
        story.append(col_tbl)
        story.append(PageBreak())

        # ── MISSING VALUES HEATMAP ────────────────────────────────────────────
        story.append(Paragraph("5. Missing Values Analysis", styles["DDH1"]))
        story.append(HRFlowable(width="100%", thickness=1,
                                color=self._rl_color(BRAND_ACCENT)))
        story.append(Spacer(1, 3*mm))
        null_bytes = ChartFactory.null_heatmap(df, figsize=(10, 3.5))
        story.append(self._img_flowable(null_bytes, 160))
        story.append(PageBreak())

        # ── DISTRIBUTIONS ─────────────────────────────────────────────────────
        story.append(Paragraph("6. Numeric Distributions", styles["DDH1"]))
        story.append(HRFlowable(width="100%", thickness=1,
                                color=self._rl_color(BRAND_ACCENT)))
        story.append(Spacer(1, 3*mm))
        dist_bytes = ChartFactory.distribution_grid(df, max_cols=6, figsize=(12, 6))
        story.append(self._img_flowable(dist_bytes, 160))

        # ── PII ALERT ─────────────────────────────────────────────────────────
        if profile.pii_detected:
            story.append(PageBreak())
            story.append(Paragraph("⚠️ 7. PII Risk Alert", styles["DDH1"]))
            story.append(HRFlowable(width="100%", thickness=2,
                                    color=self._rl_color(BRAND_DANGER)))
            story.append(Spacer(1, 3*mm))
            story.append(Paragraph(
                "The following columns may contain Personally Identifiable Information (PII). "
                "Review and apply appropriate anonymization, masking, or encryption before "
                "sharing or deploying this dataset.",
                styles["DDBody"],
            ))
            story.append(Spacer(1, 3*mm))
            pii_data = [["Column", "PII Signals Detected"]]
            for col, signals in profile.pii_detected.items():
                pii_data.append([col, ", ".join(signals)])
            pii_tbl = Table(pii_data, colWidths=[60*mm, 100*mm], repeatRows=1)
            pii_tbl.setStyle(TableStyle([
                ("BACKGROUND",   (0, 0), (-1, 0), self._rl_color(BRAND_DANGER)),
                ("TEXTCOLOR",    (0, 0), (-1, 0), rlcolors.white),
                ("FONTNAME",     (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE",     (0, 0), (-1, 0), 9),
                ("FONTNAME",     (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE",     (0, 1), (-1, -1), 9),
                ("ROWBACKGROUNDS",(0, 1), (-1, -1),
                 [rlcolors.mistyrose, rlcolors.white]),
                ("GRID",         (0, 0), (-1, -1), 0.4, self._rl_color(BRAND_DANGER)),
                ("BOTTOMPADDING",(0, 0), (-1, -1), 5),
                ("TOPPADDING",   (0, 0), (-1, -1), 5),
            ]))
            story.append(pii_tbl)

        doc.build(story, onFirstPage=self._header_footer,
                  onLaterPages=self._header_footer)
        logger.info(f"[PDF] Saved → {output_path}")
        return output_path


# ─────────────────────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
#  DOCX REPORT GENERATOR
# ══════════════════════════════════════════════════════════════════════════════
# ─────────────────────────────────────────────────────────────────────────────

class DOCXReportGenerator:
    """
    Professional Word document report using python-docx.

    Sections:
        1. Title + metadata table
        2. Executive Summary
        3. Dimension Scores table + embedded chart
        4. Issues & Recommendations
        5. Dataset Statistics
        6. Column Details table
        7. Charts (null heatmap, distributions)
        8. PII Alert
    """

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        h = hex_color.lstrip("#")
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    def _add_heading(self, doc, text: str, level: int = 1):
        from docx.shared import RGBColor, Pt
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.color.rgb = RGBColor(*self._hex_to_rgb(BRAND_PRIMARY))
        return h

    def _add_image_from_bytes(self, doc, img_bytes: bytes, width_cm: float = 15.0):
        from docx.shared import Cm
        buf = io.BytesIO(img_bytes)
        doc.add_picture(buf, width=Cm(width_cm))

    def _add_dimension_table(self, doc, dimensions: Dict[str, Any]):
        from docx.shared import RGBColor, Pt, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        headers = ["Dimension", "Score", "Weight", "Status", "Issues"]
        tbl = doc.add_table(rows=1 + len(dimensions), cols=len(headers))
        tbl.style = "Table Grid"

        # header row
        hdr_cells = tbl.rows[0].cells
        for i, h in enumerate(headers):
            cell = hdr_cells[i]
            cell.text = h
            run = cell.paragraphs[0].runs[0]
            run.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            # background
            tc   = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd  = OxmlElement("w:shd")
            shd.set(qn("w:fill"), BRAND_PRIMARY.lstrip("#"))
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:val"), "clear")
            tcPr.append(shd)

        # data rows
        for row_i, dim in enumerate(dimensions.values(), 1):
            cells = tbl.rows[row_i].cells
            status = (
                "✅ Excellent" if dim.score >= 85 else
                "🟡 Good"     if dim.score >= 70 else
                "🟠 Fair"     if dim.score >= 55 else
                "🔴 Poor"
            )
            values = [
                dim.name,
                f"{dim.score:.1f}/100",
                f"{dim.weight * 100:.0f}%",
                status,
                str(len(dim.issues)),
            ]
            for j, val in enumerate(values):
                cells[j].text = val

    def _add_stats_table(self, doc, stats: Dict[str, Any]):
        from docx.shared import RGBColor
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        labels = {
            "rows": "Total Rows", "columns": "Total Columns",
            "numeric_columns": "Numeric Columns",
            "categorical_cols": "Categorical Columns",
            "null_cells": "Null Cells", "null_rate": "Null Rate",
            "duplicate_rows": "Duplicate Rows",
            "memory_mb": "Memory (MB)",
        }
        tbl = doc.add_table(rows=1 + len(labels), cols=2)
        tbl.style = "Table Grid"

        # header
        hdr = tbl.rows[0].cells
        for cell, text in zip(hdr, ["Metric", "Value"]):
            cell.text = text
            run = cell.paragraphs[0].runs[0]
            run.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            tc   = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd  = OxmlElement("w:shd")
            shd.set(qn("w:fill"), BRAND_ACCENT.lstrip("#"))
            shd.set(qn("w:val"), "clear")
            tcPr.append(shd)

        for row_i, (key, label) in enumerate(labels.items(), 1):
            val = stats.get(key, "N/A")
            if key == "null_rate":
                val = f"{float(val):.2%}"
            tbl.rows[row_i].cells[0].text = label
            tbl.rows[row_i].cells[1].text = str(val)

    def generate(
        self,
        profile,
        df: pd.DataFrame,
        output_path: str,
        dataset_name: str = "Dataset",
    ) -> str:
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # ── page margins ──────────────────────────────────────────────────────
        for section in doc.sections:
            section.top_margin    = Cm(2.0)
            section.bottom_margin = Cm(2.0)
            section.left_margin   = Cm(2.5)
            section.right_margin  = Cm(2.5)

        # ── TITLE BLOCK ───────────────────────────────────────────────────────
        title_p = doc.add_paragraph()
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title_p.add_run("🩺 dataDoctor — Data Quality Report")
        run.bold     = True
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(*self._hex_to_rgb(BRAND_PRIMARY))

        sub_p = doc.add_paragraph()
        sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub_run = sub_p.add_run(
            f"Dataset: {dataset_name}  |  "
            f"{datetime.now().strftime('%Y-%m-%d %H:%M')}  |  "
            f"Score: {profile.overall_score:.1f}/100 — Grade {profile.grade}"
        )
        sub_run.font.size  = Pt(11)
        sub_run.font.color.rgb = RGBColor(*self._hex_to_rgb(BRAND_GRAY))

        doc.add_paragraph()

        # ── EXECUTIVE SUMMARY ─────────────────────────────────────────────────
        self._add_heading(doc, "1. Executive Summary", level=1)
        total_issues = sum(len(d.issues) for d in profile.dimensions.values())
        summary = (
            f"This report evaluates {dataset_name} across 7 quality dimensions. "
            f"The dataset ({profile.dataset_shape[0]:,} rows × {profile.dataset_shape[1]} cols) "
            f"received an overall score of {profile.overall_score:.1f}/100 "
            f"({profile.grade_label}). "
            f"{total_issues} total issues were identified."
        )
        doc.add_paragraph(summary)
        if profile.pii_detected:
            p = doc.add_paragraph()
            run = p.add_run(f"⚠ PII DETECTED in {len(profile.pii_detected)} column(s). See Section 8.")
            run.bold = True
            run.font.color.rgb = RGBColor(*self._hex_to_rgb(BRAND_DANGER))

        # score gauge chart
        self._add_heading(doc, "Overall Quality Score", level=2)
        gauge_bytes = ChartFactory.score_gauge(profile.overall_score, profile.grade)
        self._add_image_from_bytes(doc, gauge_bytes, width_cm=8)

        # ── DIMENSION SCORES ──────────────────────────────────────────────────
        self._add_heading(doc, "2. Dimension Scores", level=1)
        self._add_dimension_table(doc, profile.dimensions)
        doc.add_paragraph()

        # dimension bar chart
        bar_bytes = ChartFactory.dimension_bar_chart(
            profile.dimensions, profile.overall_score
        )
        self._add_image_from_bytes(doc, bar_bytes, width_cm=15)

        # ── PER-DIMENSION DETAIL ──────────────────────────────────────────────
        self._add_heading(doc, "3. Dimension Details", level=1)
        for dim in profile.dimensions.values():
            self._add_heading(doc, f"{dim.name} — {dim.score:.1f}/100", level=2)
            if dim.positives:
                doc.add_paragraph("Positives:", style="List Bullet")
                for p in dim.positives[:3]:
                    doc.add_paragraph(p, style="List Bullet")
            if dim.issues:
                doc.add_paragraph("Issues:", style="List Bullet")
                for issue in dim.issues[:4]:
                    para = doc.add_paragraph(style="List Bullet")
                    run  = para.add_run(f"⚠ {issue}")
                    run.font.color.rgb = RGBColor(*self._hex_to_rgb(BRAND_DANGER))

        # ── RECOMMENDATIONS ───────────────────────────────────────────────────
        self._add_heading(doc, "4. Recommendations", level=1)
        for i, rec in enumerate(profile.recommendations, 1):
            doc.add_paragraph(f"{i}. {rec}", style="List Number")

        # ── STATISTICS ───────────────────────────────────────────────────────
        self._add_heading(doc, "5. Dataset Statistics", level=1)
        self._add_stats_table(doc, profile.statistics)

        # ── MISSING VALUES CHART ──────────────────────────────────────────────
        self._add_heading(doc, "6. Missing Values", level=1)
        null_bytes = ChartFactory.null_heatmap(df, figsize=(10, 3.5))
        self._add_image_from_bytes(doc, null_bytes, width_cm=15)

        # ── DISTRIBUTIONS ─────────────────────────────────────────────────────
        self._add_heading(doc, "7. Numeric Distributions", level=1)
        dist_bytes = ChartFactory.distribution_grid(df)
        self._add_image_from_bytes(doc, dist_bytes, width_cm=15)

        # ── PII ALERT ─────────────────────────────────────────────────────────
        if profile.pii_detected:
            self._add_heading(doc, "8. ⚠ PII Risk Alert", level=1)
            p = doc.add_paragraph(
                "The following columns contain possible PII. "
                "Apply anonymization before sharing."
            )
            tbl = doc.add_table(rows=1 + len(profile.pii_detected), cols=2)
            tbl.style = "Table Grid"
            tbl.rows[0].cells[0].text = "Column"
            tbl.rows[0].cells[1].text = "Signals"
            for i, (col, signals) in enumerate(profile.pii_detected.items(), 1):
                tbl.rows[i].cells[0].text = col
                tbl.rows[i].cells[1].text = ", ".join(signals)

        doc.save(output_path)
        logger.info(f"[DOCX] Saved → {output_path}")
        return output_path


# ─────────────────────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
#  PPTX REPORT GENERATOR
# ══════════════════════════════════════════════════════════════════════════════
# ─────────────────────────────────────────────────────────────────────────────

class PPTXReportGenerator:
    """
    Professional PowerPoint presentation using python-pptx.

    Slides:
        1.  Title slide
        2.  Agenda
        3.  Overall Score + Gauge
        4.  Dimension Scores — bar chart
        5.  Dimension Scores — table
        6.  Top Issues
        7.  Recommendations
        8.  Missing Values chart
        9.  Distributions chart
        10. Dataset Statistics
        11. PII Alert (if any)
        12. Thank You / Next Steps
    """

    def _hex_to_rgb(self, hex_color: str):
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _add_img(self, slide, img_bytes: bytes, left, top, width, height=None):
        from pptx.util import Inches, Emu
        buf = io.BytesIO(img_bytes)
        if height:
            slide.shapes.add_picture(buf, left, top, width, height)
        else:
            slide.shapes.add_picture(buf, left, top, width)

    def _title_slide(self, prs, dataset_name: str, profile):
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        slide  = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        width  = prs.slide_width
        height = prs.slide_height

        # background rect
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            0, 0, width, height,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(BRAND_PRIMARY)
        bg.line.fill.background()

        # accent bar
        bar = slide.shapes.add_shape(1, 0, int(height * 0.7), width, int(height * 0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(BRAND_ACCENT)
        bar.line.fill.background()

        # title
        txb = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), width - Inches(1), Inches(1.5)
        )
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "🩺 dataDoctor"
        run.font.bold = True
        run.font.size = Pt(40)
        run.font.color.rgb = RGBColor(255, 255, 255)

        # subtitle
        txb2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8), width - Inches(1), Inches(1.2)
        )
        tf2  = txb2.text_frame
        p2   = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = "Data Quality Intelligence Report"
        run2.font.size = Pt(22)
        run2.font.color.rgb = self._hex_to_rgb(BRAND_ACCENT)

        # dataset info
        txb3 = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.0), width - Inches(1), Inches(1.5)
        )
        tf3  = txb3.text_frame
        p3   = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = (
            f"Dataset: {dataset_name}   |   "
            f"Score: {profile.overall_score:.1f}/100   |   "
            f"Grade: {profile.grade}   |   "
            f"{datetime.now().strftime('%Y-%m-%d')}"
        )
        run3.font.size = Pt(14)
        run3.font.color.rgb = RGBColor(200, 220, 240)

        slide.notes_slide.notes_text_frame.text = (
            "Welcome to the dataDoctor Data Quality Report. "
            f"This presentation covers the quality analysis of {dataset_name}."
        )

    def _content_slide(self, prs, title: str, note: str = ""):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        slide  = prs.slides.add_slide(prs.slide_layouts[6])
        width  = prs.slide_width
        height = prs.slide_height

        # header bar
        hdr = slide.shapes.add_shape(1, 0, 0, width, Inches(0.8))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = self._hex_to_rgb(BRAND_PRIMARY)
        hdr.line.fill.background()

        # title text
        txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12),
                                       width - Inches(0.6), Inches(0.6))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.bold  = True
        run.font.size  = Pt(18)
        run.font.color.rgb = RGBColor(255, 255, 255)

        # footer
        ftr = slide.shapes.add_shape(1, 0, height - Inches(0.3),
                                     width, Inches(0.3))
        ftr.fill.solid()
        ftr.fill.fore_color.rgb = self._hex_to_rgb(BRAND_GRAY)
        ftr.line.fill.background()

        txb_ftr = slide.shapes.add_textbox(
            Inches(0.2), height - Inches(0.28), width - Inches(0.4), Inches(0.25)
        )
        tf_ftr  = txb_ftr.text_frame
        p_ftr   = tf_ftr.paragraphs[0]
        p_ftr.alignment = PP_ALIGN.RIGHT
        run_ftr = p_ftr.add_run()
        run_ftr.text = f"dataDoctor  |  {datetime.now().strftime('%Y-%m-%d')}"
        run_ftr.font.size = Pt(7)
        run_ftr.font.color.rgb = RGBColor(240, 240, 240)

        if note:
            slide.notes_slide.notes_text_frame.text = note

        return slide

    def _add_text_box(self, slide, text: str, left, top, width, height,
                      font_size: int = 11, bold: bool = False,
                      color: str = BRAND_DARK, wrap: bool = True):
        from pptx.util import Pt
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = self._hex_to_rgb(color)

    def generate(
        self,
        profile,
        df: pd.DataFrame,
        output_path: str,
        dataset_name: str = "Dataset",
    ) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        W = prs.slide_width
        H = prs.slide_height

        # ── SLIDE 1: Title ────────────────────────────────────────────────────
        self._title_slide(prs, dataset_name, profile)

        # ── SLIDE 2: Agenda ───────────────────────────────────────────────────
        s2 = self._content_slide(prs, "Agenda", "Overview of today's presentation.")
        items = [
            "1. Overall Quality Score",
            "2. Dimension Breakdown",
            "3. Key Issues",
            "4. Recommendations",
            "5. Missing Values Analysis",
            "6. Distributions",
            "7. Dataset Statistics",
            "8. PII Alert" if profile.pii_detected else "8. Summary",
        ]
        for i, item in enumerate(items):
            self._add_text_box(
                s2, item,
                Inches(1.5), Inches(0.9 + i * 0.7),
                Inches(10), Inches(0.6),
                font_size=14,
            )

        # ── SLIDE 3: Score ────────────────────────────────────────────────────
        s3 = self._content_slide(
            prs, "Overall Quality Score",
            f"Score: {profile.overall_score:.1f}/100, Grade: {profile.grade}"
        )
        gauge_bytes = ChartFactory.score_gauge(
            profile.overall_score, profile.grade, figsize=(5, 5)
        )
        self._add_img(s3, gauge_bytes, Inches(0.3), Inches(0.9),
                      Inches(5), Inches(5))

        self._add_text_box(
            s3, f"Grade: {profile.grade_label}",
            Inches(5.8), Inches(1.5), Inches(7), Inches(0.8),
            font_size=22, bold=True, color=BRAND_PRIMARY,
        )
        self._add_text_box(
            s3, profile.health_badge,
            Inches(5.8), Inches(2.4), Inches(7), Inches(0.7),
            font_size=18, color=BRAND_ACCENT,
        )
        self._add_text_box(
            s3,
            f"Rows: {profile.dataset_shape[0]:,}   |   "
            f"Columns: {profile.dataset_shape[1]}   |   "
            f"Numeric: {len(profile.numeric_columns)}",
            Inches(5.8), Inches(3.2), Inches(7), Inches(0.6),
            font_size=13,
        )
        total_issues = sum(len(d.issues) for d in profile.dimensions.values())
        self._add_text_box(
            s3,
            f"Total Issues Found: {total_issues}",
            Inches(5.8), Inches(3.9), Inches(7), Inches(0.6),
            font_size=13,
            color=BRAND_DANGER if total_issues > 5 else BRAND_SUCCESS,
        )

        # ── SLIDE 4: Dimension Bar Chart ──────────────────────────────────────
        s4 = self._content_slide(prs, "Quality Dimension Scores",
                                 "Each bar shows the score for one dimension.")
        bar_bytes = ChartFactory.dimension_bar_chart(
            profile.dimensions, profile.overall_score, figsize=(11, 5)
        )
        self._add_img(s4, bar_bytes, Inches(0.3), Inches(0.85),
                      Inches(12.7), Inches(5.8))

        # ── SLIDE 5: Dimension Table ──────────────────────────────────────────
        s5 = self._content_slide(prs, "Dimension Score Summary")
        dims     = list(profile.dimensions.values())
        rows_tbl = 1 + len(dims)
        cols_tbl = 5
        col_widths = [Inches(2.8), Inches(1.8), Inches(1.5), Inches(2.5), Inches(1.5)]

        tbl = s5.shapes.add_table(
            rows_tbl, cols_tbl, Inches(0.3), Inches(0.9),
            sum(col_widths), Inches(0.4 * rows_tbl)
        ).table
        tbl.columns[0].width = col_widths[0]
        tbl.columns[1].width = col_widths[1]
        tbl.columns[2].width = col_widths[2]
        tbl.columns[3].width = col_widths[3]
        tbl.columns[4].width = col_widths[4]

        headers = ["Dimension", "Score /100", "Weight", "Status", "Issues"]
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._hex_to_rgb(BRAND_PRIMARY)

        for i, dim in enumerate(dims, 1):
            status = (
                "✅ Excellent" if dim.score >= 85 else
                "🟡 Good"     if dim.score >= 70 else
                "🟠 Fair"     if dim.score >= 55 else "🔴 Poor"
            )
            vals = [dim.name, f"{dim.score:.1f}", f"{dim.weight*100:.0f}%",
                    status, str(len(dim.issues))]
            for j, val in enumerate(vals):
                cell = tbl.cell(i, j)
                cell.text = val
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 245, 250)

        # ── SLIDE 6: Top Issues ───────────────────────────────────────────────
        s6 = self._content_slide(prs, "Top Issues Identified")
        all_issues = [i for d in profile.dimensions.values() for i in d.issues][:8]
        for idx, issue in enumerate(all_issues):
            self._add_text_box(
                s6, f"• {issue[:90]}",
                Inches(0.5), Inches(0.9 + idx * 0.75),
                Inches(12.3), Inches(0.65),
                font_size=11, color=BRAND_DANGER,
            )

        # ── SLIDE 7: Recommendations ──────────────────────────────────────────
        s7 = self._content_slide(prs, "Recommendations")
        for idx, rec in enumerate(profile.recommendations[:7]):
            self._add_text_box(
                s7, f"{idx+1}. {rec[:100]}",
                Inches(0.5), Inches(0.9 + idx * 0.8),
                Inches(12.3), Inches(0.7),
                font_size=11,
            )

        # ── SLIDE 8: Null Heatmap ─────────────────────────────────────────────
        s8 = self._content_slide(prs, "Missing Values Analysis")
        null_bytes = ChartFactory.null_heatmap(df, figsize=(12, 4))
        self._add_img(s8, null_bytes, Inches(0.3), Inches(0.9),
                      Inches(12.7), Inches(5.5))

        # ── SLIDE 9: Distributions ────────────────────────────────────────────
        s9 = self._content_slide(prs, "Numeric Distributions")
        dist_bytes = ChartFactory.distribution_grid(df, max_cols=6, figsize=(13, 5))
        self._add_img(s9, dist_bytes, Inches(0.2), Inches(0.85),
                      Inches(12.9), Inches(6))

        # ── SLIDE 10: Statistics ──────────────────────────────────────────────
        s10 = self._content_slide(prs, "Dataset Statistics")
        stats_items = [
            ("Rows",               f"{profile.statistics.get('rows', 0):,}"),
            ("Columns",            str(profile.statistics.get("columns", 0))),
            ("Numeric Columns",    str(profile.statistics.get("numeric_columns", 0))),
            ("Null Cells",         str(profile.statistics.get("null_cells", 0))),
            ("Overall Null Rate",  f"{profile.statistics.get('null_rate', 0):.2%}"),
            ("Duplicate Rows",     str(profile.statistics.get("duplicate_rows", 0))),
            ("Memory (MB)",        str(profile.statistics.get("memory_mb", 0))),
        ]
        for idx, (label, val) in enumerate(stats_items):
            col = 0 if idx < 4 else 1
            row = idx if idx < 4 else idx - 4
            self._add_text_box(
                s10, f"{label}: {val}",
                Inches(0.5 + col * 6.5), Inches(1.0 + row * 1.0),
                Inches(6), Inches(0.8),
                font_size=14, bold=True,
            )

        # ── SLIDE 11: PII Alert ───────────────────────────────────────────────
        if profile.pii_detected:
            s11 = self._content_slide(
                prs, "⚠️ PII Risk Alert",
                "Personally Identifiable Information detected. Immediate action required."
            )
            self._add_text_box(
                s11,
                f"PII detected in {len(profile.pii_detected)} column(s). "
                "Apply anonymization before deployment.",
                Inches(0.5), Inches(0.9), Inches(12), Inches(0.8),
                font_size=14, bold=True, color=BRAND_DANGER,
            )
            for idx, (col, signals) in enumerate(profile.pii_detected.items()):
                self._add_text_box(
                    s11,
                    f"• {col}: {', '.join(signals[:3])}",
                    Inches(0.8), Inches(1.8 + idx * 0.65),
                    Inches(11.5), Inches(0.6),
                    font_size=12, color=BRAND_DANGER,
                )

        # ── SLIDE 12: Thank You ───────────────────────────────────────────────
        s12 = self._content_slide(prs, "Next Steps & Summary")
        next_steps = [
            "1. Address critical issues (🔴) before deploying the dataset",
            "2. Impute or remove missing values in key columns",
            "3. Anonymize any detected PII columns",
            "4. Re-run dataDoctor after cleaning to verify improvement",
            "5. Export cleaned pipeline via CLI: python cli.py pipeline <file>",
        ]
        for idx, step in enumerate(next_steps):
            self._add_text_box(
                s12, step,
                Inches(0.5), Inches(0.9 + idx * 0.9),
                Inches(12.3), Inches(0.75),
                font_size=13,
            )
        self._add_text_box(
            s12, "🩺 dataDoctor — github.com/Denterio1/dataDoctor",
            Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
            font_size=10, color=BRAND_GRAY,
        )

        prs.save(output_path)
        logger.info(f"[PPTX] Saved → {output_path}")
        return output_path


# ─────────────────────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
#  SMART REPORT — main entry point
# ══════════════════════════════════════════════════════════════════════════════
# ─────────────────────────────────────────────────────────────────────────────

class SmartReport:
    """
    Main entry point for dataDoctor smart report generation.

    Generates PDF, DOCX, and/or PPTX from a QualityProfile.

    Usage
    -----
    >>> from src.data.quality_score import DataQualityScorer
    >>> from src.smart_report import SmartReport
    >>>
    >>> profile = DataQualityScorer().score(df)
    >>> reporter = SmartReport(output_dir="reports/")
    >>> paths = reporter.generate_all(profile, df, dataset_name="Sales_2024")
    >>> print(paths)
    {'pdf': 'reports/Sales_2024_quality.pdf',
     'docx': 'reports/Sales_2024_quality.docx',
     'pptx': 'reports/Sales_2024_quality.pptx'}
    """

    def __init__(
        self,
        output_dir: str = ".",
        brand_name: str = "dataDoctor",
    ):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.brand_name = brand_name

    def _safe_name(self, dataset_name: str) -> str:
        return "".join(c if c.isalnum() or c in "_-" else "_" for c in dataset_name)

    def generate_pdf(
        self,
        profile,
        df: pd.DataFrame,
        dataset_name: str = "Dataset",
        filename: Optional[str] = None,
    ) -> str:
        safe = self._safe_name(dataset_name)
        path = str(self.output_dir / (filename or f"{safe}_quality.pdf"))
        gen  = PDFReportGenerator(self.brand_name)
        return gen.generate(profile, df, path, dataset_name)

    def generate_docx(
        self,
        profile,
        df: pd.DataFrame,
        dataset_name: str = "Dataset",
        filename: Optional[str] = None,
    ) -> str:
        safe = self._safe_name(dataset_name)
        path = str(self.output_dir / (filename or f"{safe}_quality.docx"))
        gen  = DOCXReportGenerator()
        return gen.generate(profile, df, path, dataset_name)

    def generate_pptx(
        self,
        profile,
        df: pd.DataFrame,
        dataset_name: str = "Dataset",
        filename: Optional[str] = None,
    ) -> str:
        safe = self._safe_name(dataset_name)
        path = str(self.output_dir / (filename or f"{safe}_quality.pptx"))
        gen  = PPTXReportGenerator()
        return gen.generate(profile, df, path, dataset_name)

    def generate_all(
        self,
        profile,
        df: pd.DataFrame,
        dataset_name: str = "Dataset",
        formats: Optional[List[str]] = None,
    ) -> Dict[str, str]:
        """
        Generate reports in all requested formats.

        Parameters
        ----------
        formats : list of "pdf", "docx", "pptx". Default: all three.

        Returns
        -------
        dict mapping format → file path.
        """
        formats = formats or ["pdf", "docx", "pptx"]
        paths   = {}

        generators = {
            "pdf"  : self.generate_pdf,
            "docx" : self.generate_docx,
            "pptx" : self.generate_pptx,
        }

        for fmt in formats:
            if fmt not in generators:
                logger.warning(f"Unknown format '{fmt}'. Skipping.")
                continue
            try:
                logger.info(f"[SmartReport] Generating {fmt.upper()}…")
                paths[fmt] = generators[fmt](profile, df, dataset_name)
                logger.info(f"[SmartReport] {fmt.upper()} → {paths[fmt]}")
            except ImportError as e:
                logger.error(
                    f"[SmartReport] Missing dependency for {fmt}: {e}. "
                    f"Install with: pip install "
                    f"{'reportlab' if fmt=='pdf' else 'python-docx' if fmt=='docx' else 'python-pptx'}"
                )
            except Exception as e:
                logger.error(f"[SmartReport] {fmt.upper()} failed: {e}")

        return paths


# ─────────────────────────────────────────────────────────────────────────────
# Convenience one-liner
# ─────────────────────────────────────────────────────────────────────────────

def generate_smart_report(
    profile,
    df: pd.DataFrame,
    dataset_name: str = "Dataset",
    output_dir: str   = ".",
    formats: Optional[List[str]] = None,
) -> Dict[str, str]:
    """
    One-liner smart report generation.

    >>> paths = generate_smart_report(profile, df, "MySales", formats=["pdf","pptx"])
    """
    return SmartReport(output_dir=output_dir).generate_all(
        profile, df, dataset_name, formats=formats
    )


# ─────────────────────────────────────────────────────────────────────────────
# Smoke test
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO,
                        format="%(levelname)s | %(message)s")

    import sys
    sys.path.insert(0, str(Path(__file__).parent.parent))

    try:
        from src.data.quality_score import DataQualityScorer
    except ImportError:
        # standalone test
        import subprocess
        subprocess.run(["python", "-c", "import quality_score"])
        from quality_score import DataQualityScorer

    rng = np.random.default_rng(42)
    n   = 300

    df_test = pd.DataFrame({
        "id"       : range(n),
        "age"      : rng.integers(18, 80, n),
        "salary"   : rng.normal(50_000, 15_000, n),
        "email"    : [f"u{i}@example.com" for i in range(n)],
        "score"    : rng.uniform(0, 100, n),
        "churn"    : rng.choice([0, 1], n, p=[0.85, 0.15]),
        "category" : rng.choice(["A", "B", "C"], n),
    })
    df_test.loc[:5, "age"] = None

    profile = DataQualityScorer().score(df_test)

    reporter = SmartReport(output_dir="/tmp/datadoctor_reports")
    paths = reporter.generate_all(profile, df_test, "TestDataset",
                                  formats=["pdf", "docx", "pptx"])

    print("\n✅ Reports generated:")
    for fmt, path in paths.items():
        size = Path(path).stat().st_size / 1024 if Path(path).exists() else 0
        print(f"   {fmt.upper()}: {path}  ({size:.1f} KB)")